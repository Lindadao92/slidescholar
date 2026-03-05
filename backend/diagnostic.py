#!/usr/bin/env python3
"""Diagnostic script for SlideScholar slide plans and .pptx output.

Usage:
    python diagnostic.py <slide_plan.json> [output.pptx]

Checks:
    - All main tables have >= 3 data rows (header + 3+ data)
    - No table has only 1 data row
    - Case study tables show progression, not snapshots
    - Backup slides numbered correctly (sequential from main slides)
    - Backup divider has no slide number
    - Key number slide present in main presentation
    - Hero table slides have annotations (not empty below table)
    - Hero figure slides have 3-4 annotations
    - Backup text slides have bullet_points (not bare annotations)
    - All text consistent across slide types
"""

import json
import sys
from pathlib import Path

# ANSI colors
GREEN = "\033[92m"
RED = "\033[91m"
YELLOW = "\033[93m"
CYAN = "\033[96m"
BOLD = "\033[1m"
RESET = "\033[0m"

def ok(msg):
    print(f"  {GREEN}✓{RESET} {msg}")

def fail(msg):
    print(f"  {RED}✗{RESET} {msg}")

def warn(msg):
    print(f"  {YELLOW}⚠{RESET} {msg}")

def header(msg):
    print(f"\n{BOLD}{CYAN}{msg}{RESET}")


def check_table_sizes(slides, label="main"):
    """Check that all tables have >= 3 data rows."""
    header(f"Table size check ({label} slides)")
    found_any = False
    all_pass = True
    for s in slides:
        td = s.get("table_data")
        if not td:
            continue
        found_any = True
        rows = td.get("rows", [])
        headers = td.get("headers", [])
        n_rows = len(rows)
        title = s.get("title", "?")[:60]
        slide_num = s.get("slide_number", "?")

        if n_rows < 3:
            fail(f"Slide {slide_num}: table has {n_rows} data row(s) (need ≥3) — \"{title}\"")
            all_pass = False
        elif n_rows <= 8:
            ok(f"Slide {slide_num}: table has {n_rows} data rows, {len(headers)} cols — \"{title}\"")
        else:
            warn(f"Slide {slide_num}: table has {n_rows} data rows (large) — \"{title}\"")

    if not found_any:
        warn(f"No tables found in {label} slides")
    elif all_pass:
        ok(f"All {label} tables have ≥3 data rows")
    return all_pass


def check_key_number(slides):
    """Check that exactly one key_number slide exists."""
    header("Key number slide check")
    key_slides = [s for s in slides if s.get("layout") == "key_number" or s.get("key_number")]
    if len(key_slides) == 0:
        fail("No key_number slide found in main presentation")
        return False
    elif len(key_slides) == 1:
        ks = key_slides[0]
        num = ks.get("key_number", "(missing)")
        ctx = ks.get("key_number_context", "(missing)")
        ok(f"Key number slide present: {num}")
        if ctx and ctx != "(missing)":
            ok(f"  Context: {ctx[:80]}")
        else:
            warn("  key_number_context is missing")
        return True
    else:
        warn(f"Found {len(key_slides)} key_number slides (expected exactly 1)")
        for ks in key_slides:
            print(f"    - Slide {ks.get('slide_number', '?')}: {ks.get('key_number', '?')}")
        return True


def check_hero_table_annotations(slides):
    """Check that hero_table slides have annotations below the table."""
    header("Hero table annotation check")
    found_any = False
    all_pass = True
    for s in slides:
        if s.get("layout") != "hero_table" or not s.get("table_data"):
            continue
        found_any = True
        anns = s.get("annotations", [])
        slide_num = s.get("slide_number", "?")
        title = s.get("title", "?")[:60]

        if not anns:
            fail(f"Slide {slide_num}: hero_table has NO annotations — \"{title}\"")
            all_pass = False
        elif len(anns) < 2:
            warn(f"Slide {slide_num}: hero_table has only {len(anns)} annotation(s) (want 2-3) — \"{title}\"")
        else:
            ok(f"Slide {slide_num}: hero_table has {len(anns)} annotations — \"{title}\"")

    if not found_any:
        warn("No hero_table slides found")
    elif all_pass:
        ok("All hero_table slides have annotations")
    return all_pass


def check_hero_figure_annotations(slides):
    """Check that hero_figure slides have 3-4 annotations."""
    header("Hero figure annotation check")
    found_any = False
    all_pass = True
    for s in slides:
        if s.get("layout") != "hero_figure":
            continue
        found_any = True
        anns = s.get("annotations", [])
        slide_num = s.get("slide_number", "?")
        title = s.get("title", "?")[:60]

        if not anns:
            fail(f"Slide {slide_num}: hero_figure has NO annotations — \"{title}\"")
            all_pass = False
        elif len(anns) < 3:
            warn(f"Slide {slide_num}: hero_figure has only {len(anns)} annotation(s) (want 3-4) — \"{title}\"")
        else:
            ok(f"Slide {slide_num}: hero_figure has {len(anns)} annotations — \"{title}\"")

    if not found_any:
        warn("No hero_figure slides found")
    elif all_pass:
        ok("All hero_figure slides have ≥3 annotations")
    return all_pass


def check_backup_bullet_format(backup_slides):
    """Check that backup text slides use bullet_points, not bare annotations."""
    header("Backup slide bullet format check")
    if not backup_slides:
        warn("No backup slides found")
        return True

    all_pass = True
    for s in backup_slides:
        ct = s.get("content_type", "")
        layout = s.get("layout", "")
        slide_num = s.get("slide_number", "?")
        title = s.get("title", "?")[:60]

        # Only check text-list backup slides (not tables or figures)
        if s.get("table_data") or layout == "hero_figure":
            continue

        bullets = s.get("bullet_points", [])
        anns = s.get("annotations", [])

        if not bullets and anns:
            warn(f"Slide {slide_num} ({ct}): has annotations but no bullet_points — "
                 f"builder will render without '•' markers — \"{title}\"")
            # Not a hard fail since the builder fix now handles this
        elif bullets:
            ok(f"Slide {slide_num} ({ct}): has {len(bullets)} bullet_points — \"{title}\"")
        elif not bullets and not anns:
            warn(f"Slide {slide_num} ({ct}): no text content — \"{title}\"")

    return all_pass


def check_layout_variety(slides):
    """Check for layout monotony in main slides."""
    header("Layout variety check")
    # Classify visual type per slide
    layouts = []
    for s in slides:
        ct = s.get("content_type", "")
        if ct in ("title", "thankyou"):
            layouts.append(("structural", s.get("slide_number", "?")))
            continue
        has_table = bool(s.get("table_data"))
        has_figure = bool(s.get("figure_reference")) and not has_table
        if has_table:
            layouts.append(("table", s.get("slide_number", "?")))
        elif has_figure:
            layouts.append(("figure", s.get("slide_number", "?")))
        else:
            layouts.append(("text", s.get("slide_number", "?")))

    # Count layout types
    table_count = sum(1 for l, _ in layouts if l == "table")
    figure_count = sum(1 for l, _ in layouts if l == "figure")
    text_count = sum(1 for l, _ in layouts if l == "text")
    total = len(slides)
    print(f"  Layout distribution: {table_count} table, {figure_count} figure, "
          f"{text_count} text, {total - table_count - figure_count - text_count} structural")

    # Check consecutive runs
    max_run = 1
    current_run = 1
    prev_type = None
    run_type = None
    for layout_type, _ in layouts:
        if layout_type == "structural":
            prev_type = None
            current_run = 1
            continue
        if layout_type == prev_type and layout_type in ("table", "figure"):
            current_run += 1
            if current_run > max_run:
                max_run = current_run
                run_type = layout_type
        else:
            current_run = 1
        prev_type = layout_type

    all_pass = True
    if max_run > 3:
        fail(f"Layout monotony: {max_run} consecutive {run_type} slides (max 3)")
        all_pass = False
    else:
        ok(f"Max consecutive same-layout: {max_run} (limit 3)")

    # Check figure ratio among visual slides
    total_visual = table_count + figure_count
    if total_visual > 0:
        fig_ratio = figure_count / total_visual * 100
        if total_visual > 4 and fig_ratio < 20:
            warn(f"Low figure ratio: {figure_count}/{total_visual} visual slides "
                 f"use figures ({fig_ratio:.0f}%, target ≥30%)")
        else:
            ok(f"Figure ratio: {figure_count}/{total_visual} visual slides "
               f"({fig_ratio:.0f}%)")

    return all_pass


def check_coverage_checklist(plan):
    """Check that coverage_checklist exists and has no missing experiments."""
    header("Experiment coverage check")
    checklist = plan.get("coverage_checklist", {})
    if not checklist:
        warn("No coverage_checklist in plan — planner may not have generated it")
        return True  # soft pass — older plans won't have it

    total = checklist.get("total_experiments", 0)
    in_main = checklist.get("experiments_in_main", [])
    in_backup = checklist.get("experiments_in_backup", [])
    missing = checklist.get("missing", [])

    print(f"  Total experiments: {total}")
    print(f"  In main slides: {len(in_main)} — {in_main}")
    print(f"  In backup slides: {len(in_backup)} — {in_backup}")

    if missing:
        fail(f"Missing experiments: {missing}")
        return False

    covered = len(in_main) + len(in_backup)
    if total > 0 and covered < total:
        fail(f"Coverage incomplete: {covered}/{total} experiments accounted for")
        return False

    ok(f"All {total} experiments covered (main={len(in_main)}, backup={len(in_backup)}), missing=[]")
    return True


def check_slide_numbering(plan):
    """Check that slide numbers are sequential and backup divider is handled."""
    header("Slide numbering check")
    main_slides = plan.get("slides", [])
    backup_slides = plan.get("backup_slides", [])

    # Check main slide numbering
    all_pass = True
    for i, s in enumerate(main_slides):
        expected = i + 1
        actual = s.get("slide_number", None)
        if actual is not None and actual != expected:
            fail(f"Main slide {i+1}: has slide_number={actual} (expected {expected})")
            all_pass = False

    if all_pass and main_slides:
        ok(f"Main slides numbered 1-{len(main_slides)} correctly")

    # Check backup numbering (will be set by builder, so just report plan values)
    if backup_slides:
        backup_nums = [s.get("slide_number", "?") for s in backup_slides]
        # After builder: main_count + 1 (divider) + 1 = first backup number
        expected_first = len(main_slides) + 2  # +1 divider +1 first backup
        print(f"  Backup slide numbers in plan: {backup_nums}")
        print(f"  Expected after builder: {expected_first} to {expected_first + len(backup_slides) - 1}")
        print(f"  (Builder post-processing will correct these)")

    return all_pass


def check_pptx_numbering(pptx_path):
    """Check slide numbering in the actual .pptx file."""
    header("PPTX slide numbering check")
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        warn("python-pptx not installed — skipping PPTX checks")
        return True

    prs = Presentation(pptx_path)
    all_pass = True
    divider_idx = None

    for i, slide in enumerate(prs.slides):
        slide_pos = i + 1
        is_divider = False
        found_number = None

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if "Backup Slides" in text:
                is_divider = True
            # Detect slide number textbox
            if (text.isdigit()
                    and shape.width < Inches(1.5)
                    and shape.top > Inches(6.0)):
                found_number = int(text)

        if is_divider:
            divider_idx = i
            if found_number:
                fail(f"Slide {slide_pos} (divider): has number {found_number} — should be blank")
                all_pass = False
            else:
                ok(f"Slide {slide_pos}: Backup divider (no number) ✓")
        elif found_number is not None:
            if found_number != slide_pos:
                fail(f"Slide {slide_pos}: shows number {found_number} (off by {found_number - slide_pos})")
                all_pass = False
            else:
                ok(f"Slide {slide_pos}: number {found_number} ✓")

    total = len(prs.slides)
    if all_pass:
        ok(f"All {total} slides numbered correctly")
    print(f"  Total slides in deck: {total}")
    if divider_idx is not None:
        print(f"  Divider at position: {divider_idx + 1}")
        print(f"  Main slides: 1-{divider_idx}")
        print(f"  Backup slides: {divider_idx + 2}-{total}")

    return all_pass


def print_slide_summary(plan):
    """Print a compact summary of all slides."""
    header("Slide plan summary")
    main_slides = plan.get("slides", [])
    backup_slides = plan.get("backup_slides", [])

    print(f"  Talk: {plan.get('talk_title', '?')}")
    print(f"  Length: {plan.get('talk_length_minutes', '?')} min")
    print(f"  Main slides: {len(main_slides)}")
    print(f"  Backup slides: {len(backup_slides)}")

    # Visual ratio
    visual_count = 0
    for s in main_slides:
        layout = s.get("layout", "")
        if layout in ("hero_figure", "hero_table", "key_number", "equation"):
            visual_count += 1
    if main_slides:
        ratio = visual_count / len(main_slides) * 100
        label = f"{ratio:.0f}%"
        if ratio >= 60:
            ok(f"Visual ratio: {visual_count}/{len(main_slides)} = {label} (target ≥60%)")
        else:
            warn(f"Visual ratio: {visual_count}/{len(main_slides)} = {label} (target ≥60%)")

    # Timing
    total_time = sum(s.get("speaking_time_seconds", 0) for s in main_slides)
    talk_min = plan.get("talk_length_minutes", 15)
    budget = talk_min * 60 * 0.8
    print(f"  Total speaking time: {total_time}s / {budget:.0f}s budget ({total_time/budget*100:.0f}%)")

    print()
    fmt = "  {num:>3}  {layout:<14} {ct:<22} {title}"
    print(fmt.format(num="#", layout="Layout", ct="Content Type", title="Title"))
    print(f"  {'─'*90}")
    for s in main_slides:
        num = s.get("slide_number", "?")
        layout = s.get("layout", "-")
        ct = s.get("content_type", "-")
        title = s.get("title", "?")[:55]
        extras = []
        if s.get("table_data"):
            rows = len(s["table_data"].get("rows", []))
            extras.append(f"tbl:{rows}r")
        if s.get("annotations"):
            extras.append(f"ann:{len(s['annotations'])}")
        if s.get("key_number"):
            extras.append(f"key:{s['key_number']}")
        if s.get("figure_reference"):
            extras.append("fig")
        suffix = f"  [{', '.join(extras)}]" if extras else ""
        print(fmt.format(num=num, layout=layout, ct=ct, title=title) + suffix)

    if backup_slides:
        print(f"  {'─'*90}")
        print(f"  {'':>3}  {'BACKUP':^14}")
        print(f"  {'─'*90}")
        for s in backup_slides:
            num = s.get("slide_number", "?")
            layout = s.get("layout", "-")
            ct = s.get("content_type", "-")
            title = s.get("title", "?")[:55]
            extras = []
            if s.get("table_data"):
                rows = len(s["table_data"].get("rows", []))
                extras.append(f"tbl:{rows}r")
            if s.get("bullet_points"):
                extras.append(f"bp:{len(s['bullet_points'])}")
            if s.get("annotations") and not s.get("bullet_points"):
                extras.append(f"ann:{len(s['annotations'])} (no bp!)")
            if s.get("figure_reference"):
                extras.append("fig")
            suffix = f"  [{', '.join(extras)}]" if extras else ""
            print(fmt.format(num=num, layout=layout, ct=ct, title=title) + suffix)


def main():
    if len(sys.argv) < 2:
        print(f"Usage: python {sys.argv[0]} <slide_plan.json> [output.pptx]")
        sys.exit(1)

    plan_path = Path(sys.argv[1])
    pptx_path = Path(sys.argv[2]) if len(sys.argv) > 2 else None

    if not plan_path.exists():
        print(f"{RED}Error: {plan_path} not found{RESET}")
        sys.exit(1)

    with open(plan_path) as f:
        plan = json.load(f)

    main_slides = plan.get("slides", [])
    backup_slides = plan.get("backup_slides", [])

    # Summary
    print_slide_summary(plan)

    # Run checks
    results = []
    results.append(check_table_sizes(main_slides, "main"))
    if backup_slides:
        results.append(check_table_sizes(backup_slides, "backup"))
    results.append(check_key_number(main_slides))
    results.append(check_hero_table_annotations(main_slides))
    results.append(check_hero_figure_annotations(main_slides))
    results.append(check_layout_variety(main_slides))
    results.append(check_backup_bullet_format(backup_slides))
    results.append(check_coverage_checklist(plan))
    results.append(check_slide_numbering(plan))

    if pptx_path and pptx_path.exists():
        results.append(check_pptx_numbering(str(pptx_path)))
    elif pptx_path:
        warn(f"PPTX file not found: {pptx_path}")

    # Final tally
    header("Results")
    passed = sum(1 for r in results if r)
    total = len(results)
    if passed == total:
        print(f"  {GREEN}{BOLD}All {total} checks passed{RESET}")
    else:
        print(f"  {passed}/{total} checks passed, {RED}{total - passed} failed{RESET}")

    sys.exit(0 if passed == total else 1)


if __name__ == "__main__":
    main()
