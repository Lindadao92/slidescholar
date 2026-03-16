"""Builds .pptx presentations from slide plans using python-pptx."""

import logging
import os
import re
import tempfile

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402

from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

log = logging.getLogger("slidescholar")

# Cache for image aspect ratios (path -> height/width)
_image_ratio_cache: dict[str, float] = {}

# --- Design constants ---
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Unified color palette — single source of truth per role
COLORS = {
    "title": RGBColor(0x2D, 0x2D, 0x2D),        # Slide titles
    "body": RGBColor(0x4A, 0x4A, 0x4A),          # Body text, annotations, bullets
    "accent": RGBColor(0x1E, 0x3A, 0x5F),        # Accent bar, key numbers, headlines
    "muted": RGBColor(0x99, 0x99, 0x99),          # References, captions, context lines
    "white": RGBColor(0xFF, 0xFF, 0xFF),          # Backgrounds, header text
    "table_header_bg": RGBColor(0x1E, 0x3A, 0x5F),  # Table header background (= accent)
    "table_row_alt": RGBColor(0xF5, 0xF7, 0xFA),    # Alternating table row bg
    "table_bold_bg": RGBColor(0xE8, 0xEE, 0xF5),    # Bold/highlighted row bg
}

# Backward-compatible aliases (used throughout; will be phased out)
COLOR_DARK = COLORS["title"]
COLOR_BODY = COLORS["body"]
COLOR_ACCENT = COLORS["accent"]
COLOR_LIGHT_GRAY = COLORS["muted"]
COLOR_WHITE = COLORS["white"]
COLOR_TABLE_HEADER_BG = COLORS["table_header_bg"]
COLOR_TABLE_ROW_ALT = COLORS["table_row_alt"]
COLOR_TABLE_BOLD_BG = COLORS["table_bold_bg"]

FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"
TITLE_SIZE = Pt(24)  # Assertion-evidence: smaller for longer sentence titles
BODY_SIZE = Pt(18)
SMALL_SIZE = Pt(12)
NOTES_SIZE = Pt(14)
TABLE_HEADER_SIZE = Pt(14)
TABLE_CELL_SIZE = Pt(13)
EQUATION_SIZE = Pt(16)
ANNOTATION_SIZE = Pt(15)
REFERENCE_SIZE = Pt(10)
CONTEXT_SIZE = Pt(12)
KEY_NUMBER_SIZE = Pt(72)
KEY_NUMBER_CONTEXT_SIZE = Pt(24)

# Layout margins
MARGIN_LEFT = Inches(0.8)
MARGIN_TOP = Inches(1.0)
MARGIN_RIGHT = Inches(0.8)
CONTENT_TOP = Inches(1.5)

# Reserved bottom zones — work upward from slide bottom
# Slide number sits at the very bottom-right
SLIDE_NUMBER_TOP = SLIDE_HEIGHT - Inches(0.35)     # 7.15"
# References sit above the slide number
REFERENCE_TOP = SLIDE_NUMBER_TOP - Inches(0.30)     # 6.85"
# Content (figures, tables, annotations) must end above this line
MAX_CONTENT_BOTTOM = REFERENCE_TOP - Inches(0.15)   # 6.70"

# Accent bar
BAR_HEIGHT = Inches(0.06)  # ~4.3pt — visible accent bar

# Max characters per bullet before truncation
MAX_BULLET_CHARS = 90


# =============================================================
# VERTICAL CURSOR — prevents text/image overlap
# =============================================================


class VerticalCursor:
    """Tracks the current Y position on a slide to prevent overlapping elements.

    Usage:
        cursor = VerticalCursor(start=CONTENT_TOP)
        cursor.place(Inches(1.0))   # returns top, advances by 1.0"
        cursor.advance(Inches(0.2)) # add gap
    """

    def __init__(self, start=0):
        self.y = int(start)

    def place(self, height, gap=0):
        """Reserve `height` EMU at the current position. Returns the top Y.

        Args:
            height: Height of the element in EMU.
            gap: Extra gap in EMU added after the element.
        """
        top = self.y
        self.y = top + int(height) + int(gap)
        return top

    def advance(self, gap):
        """Add vertical gap without placing an element."""
        self.y += int(gap)

    def remaining(self, bottom=None):
        """Return EMU remaining between current Y and the bottom boundary."""
        bottom = bottom or int(SLIDE_HEIGHT)
        return max(0, int(bottom) - self.y)

    def place_at(self, y, height, gap=0):
        """Place at a specific Y, update cursor to y + height + gap."""
        self.y = int(y) + int(height) + int(gap)
        return int(y)


# =============================================================
# SHARED HELPERS
# =============================================================


def _set_slide_bg_white(slide):
    """Set the slide background to solid white."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE


def _add_accent_bar(slide):
    """Add a thin colored bar across the top of the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left=0, top=0,
        width=SLIDE_WIDTH, height=BAR_HEIGHT,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT
    shape.line.fill.background()


def _add_slide_number(slide, number):
    """Add a slide number in the bottom-right corner (reserved zone)."""
    txBox = slide.shapes.add_textbox(
        left=SLIDE_WIDTH - Inches(1.0),
        top=SLIDE_NUMBER_TOP,
        width=Inches(0.6), height=Inches(0.3),
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(number)
    run.font.name = FONT_BODY
    run.font.size = SMALL_SIZE
    run.font.color.rgb = COLOR_ACCENT


def _add_speaker_notes(slide, slide_data):
    """Write speaker notes, transition, and timing cue to slide notes."""
    notes = slide_data.get("speaker_notes", "")
    transition = slide_data.get("transition", "")
    timing = slide_data.get("timing_cue", "")
    if not notes and not transition:
        return

    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.clear()

    if timing:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"[{timing}]"
        run.font.name = FONT_BODY
        run.font.size = NOTES_SIZE
        run.font.bold = True
        run.font.color.rgb = COLOR_ACCENT
        if notes:
            p = tf.add_paragraph()

    if notes:
        target_p = tf.paragraphs[0] if not timing else p
        if timing:
            target_p = tf.add_paragraph()
        else:
            target_p = tf.paragraphs[0]
        run = target_p.add_run()
        run.text = notes
        run.font.name = FONT_BODY
        run.font.size = NOTES_SIZE

    if transition:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"\n[Transition] {transition}"
        run.font.name = FONT_BODY
        run.font.size = NOTES_SIZE
        run.font.italic = True
        run.font.color.rgb = COLOR_LIGHT_GRAY


def _add_title_text(slide, title, left=None, top=None, width=None, size=None, alignment=None):
    """Add a styled title textbox to a slide.

    Assertion-evidence titles are full sentences (8-15 words) so the textbox
    is tall enough for 2-line wrap at the default 24pt size.
    """
    left = left if left is not None else MARGIN_LEFT
    top = top if top is not None else Inches(0.3)
    width = width if width is not None else (SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT)
    size = size or TITLE_SIZE

    # Taller box for assertion-evidence sentence titles (2–3 line wrap)
    height = Inches(1.1) if size <= Pt(28) else Inches(1.0)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = alignment or PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = FONT_TITLE
    run.font.size = size
    run.font.bold = True
    run.font.color.rgb = COLOR_DARK
    return txBox


def _truncate(text, max_chars=MAX_BULLET_CHARS):
    """Truncate text to max_chars, adding ellipsis if needed."""
    if not isinstance(text, str):
        text = str(text)
    if len(text) <= max_chars:
        return text
    return text[: max_chars - 1].rstrip() + "…"


def _add_bullets(slide, bullets, left, top, width, height=None):
    """Add a list of bullet points as a textbox."""
    height = height or Inches(4.5)
    bullets = bullets[:7]  # cap at 7 bullets to prevent overflow
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    for i, bullet in enumerate(bullets):
        if not isinstance(bullet, str):
            bullet = str(bullet)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        p.level = 0
        run = p.add_run()
        run.text = f"•  {_truncate(bullet)}"
        run.font.name = FONT_BODY
        run.font.size = BODY_SIZE
        run.font.color.rgb = COLOR_BODY

    return txBox


def _add_annotations(slide, annotations, left, top, width, height=None):
    """Add annotation labels (no bullet symbols) for assertion-evidence slides.

    Annotations are short 5-10 word technical phrases rendered in gray without
    bullet symbols, used on hero_figure, hero_table, equation, and key_number
    layouts to describe what the audience should notice.
    """
    if not annotations:
        return None
    height = height or Inches(2.5)
    txBox = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    for i, ann in enumerate(annotations[:4]):  # max 4 annotations
        if not isinstance(ann, str):
            ann = str(ann)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = _truncate(ann, 80)
        run.font.name = FONT_BODY
        run.font.size = ANNOTATION_SIZE
        run.font.color.rgb = COLOR_BODY

    return txBox


def _add_references(slide, references):
    """Add small reference citations at the bottom-left of the slide.

    Rendered at 10pt light gray, e.g. "Vaswani et al., 2017  |  Table 3 in paper"
    """
    if not references:
        return None
    ref_text = "  |  ".join(str(r) for r in references[:3])
    txBox = slide.shapes.add_textbox(
        MARGIN_LEFT, REFERENCE_TOP,
        Inches(8.0), Inches(0.25),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = ref_text
    run.font.name = FONT_BODY
    run.font.size = REFERENCE_SIZE
    run.font.color.rgb = COLOR_LIGHT_GRAY
    return txBox


def _add_context_line(slide, text, left, top, width):
    """Add a context line below a figure or table (italic, gray, small).

    Shows experimental setup info like dataset, hardware, training duration.
    """
    if not text:
        return None
    txBox = slide.shapes.add_textbox(int(left), int(top), int(width), Inches(0.35))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = _truncate(str(text), 150)
    run.font.name = FONT_BODY
    run.font.size = CONTEXT_SIZE
    run.font.italic = True
    run.font.color.rgb = COLOR_LIGHT_GRAY
    return txBox


def _add_table_headline(slide, text, left, top, width):
    """Add a bold headline sentence above a table summarizing its key finding."""
    if not text:
        return None
    txBox = slide.shapes.add_textbox(int(left), int(top), int(width), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = _truncate(str(text), 120)
    run.font.name = FONT_BODY
    run.font.size = BODY_SIZE
    run.font.bold = True
    run.font.color.rgb = COLOR_ACCENT
    return txBox


def _find_figure(figure_ref, figures):
    """Match a figure reference like 'Figure 3 (page 13)' to an extracted figure.

    Strategy (in priority order):
    1. figure_number field (set by parser's tag_figures_with_numbers)
    2. Page hint from reference — "Figure 3 (page 13)" → find image on page 13
    3. Caption match — look for "Figure 3" in the extracted caption text
    4. Index fallback — use figure number as 1-based index
    """
    if not figure_ref or not figures or not isinstance(figure_ref, str):
        return None

    ref = figure_ref.strip()

    # Extract the target figure number from the reference
    fig_num_match = re.search(r"Figure\s*(\d+)", ref, re.IGNORECASE)
    if not fig_num_match:
        return None
    target_num = int(fig_num_match.group(1))

    # Strategy 1: Direct figure_number match (most reliable)
    for fig in figures:
        if fig.get("figure_number") == target_num and os.path.isfile(fig["path"]):
            return fig

    # Strategy 2: Page hint from reference — "Figure 3 (page 13)"
    page_match = re.search(r"\(page\s*(\d+)\)", ref)
    if page_match:
        target_page = int(page_match.group(1))
        for fig in figures:
            fig_page = fig.get("page", -1)
            if (fig_page == target_page or fig_page == target_page - 1) and os.path.isfile(fig["path"]):
                return fig

    # Strategy 3: Caption match
    for fig in figures:
        caption = fig.get("caption", "")
        if re.search(rf"(?:Figure|Fig\.?)\s*{target_num}\b", caption, re.IGNORECASE):
            if os.path.isfile(fig["path"]):
                return fig

    # Strategy 4: Index fallback (1-based)
    idx = target_num - 1
    if 0 <= idx < len(figures) and os.path.isfile(figures[idx]["path"]):
        return figures[idx]

    return None


def _get_image_ratio(image_path: str) -> float:
    """Return height/width ratio of an image using PIL. Falls back to 0.75."""
    if image_path in _image_ratio_cache:
        return _image_ratio_cache[image_path]
    try:
        with PILImage.open(image_path) as img:
            w, h = img.size
            if w > 0:
                ratio = h / w
                _image_ratio_cache[image_path] = ratio
                return ratio
    except Exception:
        pass
    return 0.75  # default 4:3


def _add_image(slide, image_path, left, top, max_width, max_height,
               slide_num="?"):
    """Add an image scaled to fit within max_width×max_height, preserving aspect ratio.

    NEVER distorts. The image is scaled to the largest size that fits within
    the bounding box while maintaining the exact aspect ratio.

    Returns:
        A tuple (picture_shape, actual_width_emu, actual_height_emu) or (None, 0, 0).
    """
    try:
        img_ratio = _get_image_ratio(image_path)  # height / width

        # Scale to fit: try width-first, then constrain by height
        fig_width = int(max_width)
        fig_height = int(fig_width * img_ratio)

        if fig_height > int(max_height):
            fig_height = int(max_height)
            fig_width = int(fig_height / img_ratio)

        # Center within the bounding box
        actual_left = int(left) + (int(max_width) - fig_width) // 2
        actual_top = int(top) + (int(max_height) - fig_height) // 2

        pic = slide.shapes.add_picture(
            image_path, actual_left, actual_top, fig_width, fig_height,
        )

        # Debug logging
        fig_w_in = fig_width / 914400
        fig_h_in = fig_height / 914400
        pct = (fig_w_in * fig_h_in) / (13.333 * 7.5) * 100
        log.info("  Figure on slide %s: %.1f×%.1f in = %.0f%% of slide (ratio=%.2f)",
                 slide_num, fig_w_in, fig_h_in, pct, img_ratio)
        if pct < 25:
            log.warning("  Figure is under 25%% of slide area on slide %s", slide_num)

        return pic, fig_width, fig_height
    except Exception as exc:
        log.warning("Failed to add image %s: %s", image_path, exc)
        return None, 0, 0


def _figure_bounds_for_hero(image_path: str):
    """Compute max_width, max_height, left, top for a hero figure based on aspect ratio.

    Returns (left, top, max_width, max_height) in EMU.

    Portrait (aspect > 1.2):  right-side placement, max 5.0"×5.2"
    Landscape (aspect < 0.8): centered, max 10.0"×4.5"
    Square (0.8-1.2):         centered, max 7.0"×5.0"
    """
    ratio = _get_image_ratio(image_path)  # h/w

    if ratio > 1.2:
        # Portrait — place right side to leave room for annotations on the left
        max_w = Inches(5.0)
        max_h = Inches(5.2)
        fig_left = SLIDE_WIDTH - max_w - MARGIN_RIGHT
        fig_top = CONTENT_TOP
    elif ratio < 0.8:
        # Landscape — center horizontally
        max_w = Inches(10.0)
        max_h = Inches(4.5)
        fig_left = (SLIDE_WIDTH - max_w) // 2
        fig_top = CONTENT_TOP
    else:
        # Square-ish — center horizontally
        max_w = Inches(7.0)
        max_h = Inches(5.0)
        fig_left = (SLIDE_WIDTH - max_w) // 2
        fig_top = CONTENT_TOP

    return int(fig_left), int(fig_top), int(max_w), int(max_h)


def _add_caption(slide, text, left, top, width):
    """Add a small italic caption below a figure or table."""
    if not text:
        return
    txBox = slide.shapes.add_textbox(int(left), int(top), int(width), Inches(0.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = _truncate(text, 120)
    run.font.name = FONT_BODY
    run.font.size = SMALL_SIZE
    run.font.italic = True
    run.font.color.rgb = COLOR_LIGHT_GRAY


def _add_table(slide, table_data, left, top, width, max_height=None,
               place_caption=True):
    """Add a styled data table from table_data dict.

    Args:
        place_caption: If True (default), places caption below the table.
            Set to False when the caller manages caption/context via cursor.

    Returns:
        (table_shape, actual_table_height_emu) or (None, 0).

    table_data: {
        "headers": [...],
        "rows": [{"cells": [...], "bold": bool}],
        "caption": str,
        "highlight_terms": ["Transformer", ...]  # auto-bold cells containing these
    }
    """
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])
    if not headers or not rows:
        return None, 0

    highlight_terms = [t.lower() for t in table_data.get("highlight_terms", [])]
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header
    max_height = max_height or Inches(4.0)

    row_height = min(Inches(0.45), max_height // num_rows)
    table_height = row_height * num_rows

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        int(left), int(top), int(width), int(table_height),
    )
    table = table_shape.table

    # Style header row
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = header
        run.font.name = FONT_BODY
        run.font.size = TABLE_HEADER_SIZE
        run.font.bold = True
        run.font.color.rgb = COLOR_WHITE
        # Header background
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = COLOR_TABLE_HEADER_BG
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Style data rows
    for row_idx, row_data in enumerate(rows):
        # Handle both {"cells": [...], "bold": bool} and plain [...] formats
        if isinstance(row_data, list):
            cells = row_data
            is_bold = False
        else:
            cells = row_data.get("cells", [])
            is_bold = row_data.get("bold", False)

        for col_idx in range(num_cols):
            cell = table.cell(row_idx + 1, col_idx)
            cell_text = cells[col_idx] if col_idx < len(cells) else ""
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT

            # Auto-bold: match cell text against highlight_terms
            cell_bold = is_bold
            if not cell_bold and highlight_terms:
                cell_lower = cell_text.lower()
                cell_bold = any(term in cell_lower for term in highlight_terms)

            run = p.add_run()
            run.text = cell_text
            run.font.name = FONT_BODY
            run.font.size = TABLE_CELL_SIZE
            run.font.bold = cell_bold
            run.font.color.rgb = COLOR_DARK if cell_bold else COLOR_BODY

            # Alternating row backgrounds + bold row highlight
            cell_fill = cell.fill
            cell_fill.solid()
            if cell_bold:
                cell_fill.fore_color.rgb = COLOR_TABLE_BOLD_BG
            elif row_idx % 2 == 1:
                cell_fill.fore_color.rgb = COLOR_TABLE_ROW_ALT
            else:
                cell_fill.fore_color.rgb = COLOR_WHITE

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Caption below table (only if caller didn't opt out)
    caption = table_data.get("caption", "")
    if caption and place_caption:
        cap_top = int(top) + int(table_height) + Inches(0.15)
        _add_caption(slide, caption, left, cap_top, width)

    return table_shape, int(table_height)


def _render_equation_image(latex_str: str, output_path: str, fontsize: int = 28) -> str | None:
    """Render a LaTeX equation to a transparent PNG via matplotlib."""
    try:
        clean = latex_str.strip().strip("$").strip()
        fig, ax = plt.subplots(figsize=(10, 1.5))
        ax.text(
            0.5, 0.5, f"${clean}$",
            fontsize=fontsize, ha="center", va="center",
            transform=ax.transAxes,
        )
        ax.axis("off")
        fig.patch.set_alpha(0.0)
        fig.savefig(
            output_path, dpi=200, bbox_inches="tight",
            transparent=True, pad_inches=0.1,
        )
        plt.close(fig)
        return output_path
    except Exception as exc:
        log.debug("matplotlib equation render failed: %s", exc)
        return None


def _add_equation_text_fallback(slide, latex_str, left, top, width):
    """Fallback: render an equation as styled plain text (Unicode approximation)."""
    clean = latex_str.strip().strip("$").strip()
    clean = re.sub(r"\\text\{([^}]*)\}", r"\1", clean)
    clean = re.sub(r"\\mathrm\{([^}]*)\}", r"\1", clean)
    clean = re.sub(r"\\frac\{([^}]*)\}\{([^}]*)\}", r"(\1) / (\2)", clean)
    clean = re.sub(r"\\sqrt\{([^}]*)\}", r"√(\1)", clean)
    clean = clean.replace("^T", "ᵀ").replace("^{T}", "ᵀ")
    clean = clean.replace("^2", "²").replace("^{2}", "²")
    clean = clean.replace("_k", "ₖ").replace("_{k}", "ₖ")
    clean = clean.replace("_i", "ᵢ").replace("_{i}", "ᵢ")
    clean = clean.replace("_v", "ᵥ").replace("_{v}", "ᵥ")
    clean = clean.replace("\\cdot", "·").replace("\\times", "×")
    clean = clean.replace("\\leq", "≤").replace("\\geq", "≥")
    clean = clean.replace("\\infty", "∞").replace("\\sum", "Σ")
    clean = clean.replace("\\", "")

    txBox = slide.shapes.add_textbox(int(left), int(top), int(width), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = clean
    run.font.name = "Consolas"
    run.font.size = EQUATION_SIZE
    run.font.color.rgb = COLOR_ACCENT
    run.font.bold = True
    return txBox


def _add_equation(slide, latex_str, left, top, width):
    """Add an equation — rendered as an image via matplotlib, with text fallback."""
    if not latex_str:
        return None

    # Try matplotlib image rendering first
    eq_dir = tempfile.gettempdir()
    eq_hash = abs(hash(latex_str)) % 10**8
    eq_path = os.path.join(eq_dir, f"slidescholar_eq_{eq_hash}.png")
    rendered = _render_equation_image(latex_str, eq_path)

    if rendered and os.path.isfile(rendered):
        eq_max_h = Inches(1.2)
        eq_min_w = int(SLIDE_WIDTH * 0.40)
        try:
            img_ratio = _get_image_ratio(rendered)
            actual_w = max(int(width * 0.6), eq_min_w)
            actual_h = int(actual_w * img_ratio)
            if actual_h > int(eq_max_h):
                actual_h = int(eq_max_h)
                actual_w = int(actual_h / img_ratio)
            if actual_w < eq_min_w:
                actual_w = eq_min_w
            centered_left = int(left) + (int(width) - actual_w) // 2
            pic = slide.shapes.add_picture(
                rendered, centered_left, int(top), actual_w, actual_h,
            )
            return pic
        except Exception:
            pass  # fall through to text

    return _add_equation_text_fallback(slide, latex_str, left, top, width)


# =============================================================
# SLIDE LAYOUT FUNCTIONS
# =============================================================


def add_title_slide(prs, slide_data, paper_meta=None):
    """Title slide: centered title (once), hook, authors, venue. No duplicates."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg_white(slide)
    _add_accent_bar(slide)

    paper_meta = paper_meta or {}
    cursor = VerticalCursor(start=Inches(1.8))

    # Title — large, centered (rendered ONCE only)
    title_text = slide_data.get("title", paper_meta.get("talk_title", ""))
    title_h = Inches(1.0)
    title_top = cursor.place(title_h, gap=Inches(0.15))
    _add_title_text(
        slide, title_text,
        left=Inches(1.5), top=title_top,
        width=SLIDE_WIDTH - Inches(3.0),
        size=Pt(36), alignment=PP_ALIGN.CENTER,
    )

    # Subtitle / hook (must differ from title)
    subtitle = paper_meta.get("talk_subtitle", "")
    if subtitle and subtitle.lower().strip() != title_text.lower().strip():
        sub_h = Inches(0.5)
        sub_top = cursor.place(sub_h, gap=Inches(0.15))
        txBox = slide.shapes.add_textbox(
            Inches(2.0), sub_top,
            SLIDE_WIDTH - Inches(4.0), sub_h,
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = subtitle
        run.font.name = FONT_BODY
        run.font.size = Pt(18)
        run.font.italic = True
        run.font.color.rgb = COLOR_ACCENT

    # Authors
    authors = paper_meta.get("authors", "")
    if authors and authors != "Unknown":
        auth_h = Inches(0.5)
        auth_top = cursor.place(auth_h, gap=Inches(0.15))
        txBox = slide.shapes.add_textbox(
            Inches(1.5), auth_top,
            SLIDE_WIDTH - Inches(3.0), auth_h,
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = authors
        run.font.name = FONT_BODY
        run.font.size = Pt(20)
        run.font.color.rgb = COLOR_BODY

    # Venue
    venue = paper_meta.get("venue", "")
    if venue:
        ven_h = Inches(0.4)
        ven_top = cursor.place(ven_h)
        txBox = slide.shapes.add_textbox(
            Inches(2.0), ven_top,
            SLIDE_WIDTH - Inches(4.0), ven_h,
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = venue
        run.font.name = FONT_BODY
        run.font.size = Pt(16)
        run.font.color.rgb = COLOR_LIGHT_GRAY

    # NOTE: bullet_points are intentionally NOT rendered on the title slide.
    # The subtitle, authors, and venue are rendered from paper_meta above.
    # Claude sometimes duplicates these into bullet_points — skip them.

    _add_slide_number(slide, slide_data.get("slide_number", 1))
    _add_speaker_notes(slide, slide_data)
    return slide


def add_content_slide(prs, slide_data, figures=None):
    """Content/bullets slide: assertion title + bullets left, optional figure right.

    When a figure is present the layout splits into:
      - Left 48% for bullet/annotation text
      - Right 45% for the figure (5% gap)
    Figure height scales dynamically from the image's aspect ratio,
    capped at 70% of slide height.

    If annotations are present and no bullet_points, uses annotations.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg_white(slide)
    _add_accent_bar(slide)

    _add_title_text(slide, slide_data.get("title", ""))

    # Prefer bullet_points for "bullets" layout; fall back to annotations
    bullets = slide_data.get("bullet_points") or []
    annotations = slide_data.get("annotations") or []
    text_items = bullets if bullets else annotations
    # Backup text slides should always use bullet formatting for consistency
    is_backup_text = slide_data.get("content_type", "").startswith("backup_")
    use_annotations = bool(not bullets and annotations and not is_backup_text)

    fig_ref = slide_data.get("figure_reference")
    fig = _find_figure(fig_ref, figures) if figures else None
    equation = slide_data.get("equation_latex")

    if fig:
        # Two-column: text left 48 %, figure right 45 %
        text_width = int(SLIDE_WIDTH * 0.48)
        fig_area_w = int(SLIDE_WIDTH * 0.45)
        fig_left = int(SLIDE_WIDTH * 0.52)
        fig_max_h = int(SLIDE_HEIGHT * 0.70)
        fig_top = CONTENT_TOP

        text_h = Inches(3.5) if equation else Inches(4.5)
        if use_annotations:
            _add_annotations(slide, text_items, MARGIN_LEFT, CONTENT_TOP, text_width, text_h)
        else:
            _add_bullets(slide, text_items, MARGIN_LEFT, CONTENT_TOP, text_width, text_h)

        if equation:
            eq_top = CONTENT_TOP + text_h + Inches(0.2)
            _add_equation(slide, equation, MARGIN_LEFT, int(eq_top), text_width)

        pic, actual_w, actual_h = _add_image(
            slide, fig["path"], fig_left, int(fig_top),
            fig_area_w, fig_max_h,
            slide_num=slide_data.get("slide_number", "?"),
        )

        caption = slide_data.get("figure_caption") or fig.get("caption", "")
        if pic and caption:
            cap_top = int(fig_top) + (actual_h if actual_h else fig_max_h) + Inches(0.1)
            if cap_top + Inches(0.4) <= int(MAX_CONTENT_BOTTOM):
                _add_caption(slide, caption, fig_left, cap_top, fig_area_w)
    else:
        # Full-width text
        text_width = SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT
        text_h = Inches(3.5) if equation else Inches(4.5)
        if use_annotations:
            _add_annotations(slide, text_items, MARGIN_LEFT, CONTENT_TOP, text_width, text_h)
        else:
            _add_bullets(slide, text_items, MARGIN_LEFT, CONTENT_TOP, text_width, text_h)

        if equation:
            eq_top = CONTENT_TOP + text_h + Inches(0.2)
            _add_equation(slide, equation, Inches(2.0), int(eq_top), SLIDE_WIDTH - Inches(4.0))

    _add_references(slide, slide_data.get("references"))
    _add_slide_number(slide, slide_data.get("slide_number", ""))
    _add_speaker_notes(slide, slide_data)
    return slide


def add_table_slide(prs, slide_data, figures=None):
    """Hero table slide: assertion title + table_headline + data table + context_line.

    Uses VerticalCursor to prevent overlap between caption, context_line,
    and bottom-zone elements. Respects MAX_CONTENT_BOTTOM.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg_white(slide)
    _add_accent_bar(slide)

    _add_title_text(slide, slide_data.get("title", ""))

    table_data = slide_data.get("table_data")
    headline = slide_data.get("table_headline")
    bullets = slide_data.get("bullet_points", [])

    cursor = VerticalCursor(start=CONTENT_TOP)
    table_width = SLIDE_WIDTH - Inches(2.0)
    table_left = Inches(1.0)

    # Prefer table_headline over bullets; fall back to short bullets
    if headline:
        hl_top = cursor.place(Inches(0.4), gap=Inches(0.1))
        _add_table_headline(slide, headline, table_left, hl_top, table_width)
    elif bullets:
        bullet_h = Inches(1.2)
        bl_top = cursor.place(bullet_h, gap=Inches(0.2))
        _add_bullets(
            slide, bullets, MARGIN_LEFT, bl_top,
            SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT, bullet_h,
        )

    # The table — caller handles caption/context, so place_caption=False
    if table_data:
        remaining_h = int(MAX_CONTENT_BOTTOM) - cursor.y - Inches(1.0)
        if remaining_h < Inches(1.0):
            remaining_h = Inches(1.0)
        tbl_top = cursor.y
        table_shape, actual_table_h = _add_table(
            slide, table_data, table_left, tbl_top, table_width,
            remaining_h, place_caption=False,
        )
        cursor.place_at(tbl_top, actual_table_h, gap=Inches(0.12))

        # Caption below the table (from table_data)
        # _add_caption creates Inches(0.4) tall box
        caption = table_data.get("caption", "")
        if caption and cursor.y + Inches(0.40) <= int(MAX_CONTENT_BOTTOM):
            cap_top = cursor.place(Inches(0.40))
            _add_caption(slide, caption, table_left, cap_top, table_width)

        # Context line below caption (from slide_data)
        # _add_context_line creates Inches(0.35) tall box
        context = slide_data.get("context_line")
        if context and cursor.y + Inches(0.35) <= int(MAX_CONTENT_BOTTOM):
            ctx_top = cursor.place(Inches(0.35))
            _add_context_line(slide, context, table_left, ctx_top, table_width)

        # Annotations below context — fill remaining space with takeaways
        annotations = slide_data.get("annotations", [])
        if annotations:
            cursor.advance(Inches(0.15))
            ann_count = min(len(annotations), 3)
            ann_h = Inches(ann_count * 0.38)
            if cursor.y + ann_h <= int(MAX_CONTENT_BOTTOM):
                ann_top = cursor.place(ann_h)
                _add_annotations(
                    slide, annotations, table_left, ann_top,
                    table_width, ann_h,
                )
    else:
        if not bullets and not headline:
            bl_top = cursor.place(Inches(2.0))
            _add_bullets(
                slide, ["(Table data not available)"], MARGIN_LEFT, bl_top,
                SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
            )

    _add_references(slide, slide_data.get("references"))
    _add_slide_number(slide, slide_data.get("slide_number", ""))
    _add_speaker_notes(slide, slide_data)
    return slide


def _place_annotations_individually(slide, annotations, left, top, width, bottom_limit):
    """Place each annotation as its own textbox, advancing Y after each.

    Unlike _add_annotations (single textbox, paragraphs), this creates
    one textbox per annotation with explicit Y tracking. This prevents
    stacking/clipping issues where annotations overlap.

    Returns the Y position after the last annotation placed.
    """
    ann_y = int(top)
    ann_h = Inches(0.38)
    ann_gap = Inches(0.40)

    for ann in annotations[:4]:
        if ann_y + int(ann_h) > int(bottom_limit):
            break
        if not isinstance(ann, str):
            ann = str(ann)
        txBox = slide.shapes.add_textbox(int(left), ann_y, int(width), int(ann_h))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = _truncate(ann, 80)
        run.font.name = FONT_BODY
        run.font.size = ANNOTATION_SIZE
        run.font.color.rgb = COLOR_BODY
        ann_y += int(ann_gap)

    return ann_y


def add_figure_slide(prs, slide_data, figures=None):
    """Hero figure slide: assertion title + large figure + annotations.

    Position-based layout (same logic for main AND backup):
      - Figure on RIGHT half → annotations on LEFT side
      - Figure centered/LEFT → annotations BELOW the figure
    All elements respect MAX_CONTENT_BOTTOM.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg_white(slide)
    _add_accent_bar(slide)

    _add_title_text(slide, slide_data.get("title", ""))

    fig_ref = slide_data.get("figure_reference")
    fig = _find_figure(fig_ref, figures) if figures else None
    annotations = slide_data.get("annotations", [])

    if fig:
        # Compute figure bounds
        fig_left, fig_top, fig_max_w, fig_max_h = _figure_bounds_for_hero(fig["path"])
        available_h = int(MAX_CONTENT_BOTTOM) - int(CONTENT_TOP)
        if fig_max_h > available_h:
            fig_max_h = available_h

        # Determine if figure will be on the RIGHT half of the slide.
        # _figure_bounds_for_hero places portrait images on the right.
        fig_on_right = int(fig_left) > int(SLIDE_WIDTH * 0.35)

        # For centered figures with annotations, reserve vertical space
        # below for caption + annotations so they aren't silently dropped.
        if not fig_on_right and annotations:
            ann_count = min(len(annotations), 4)
            ann_reserve = Inches(0.60) + Inches(ann_count * 0.40)
            max_fig_h = available_h - int(ann_reserve)
            if max_fig_h < Inches(1.5):
                max_fig_h = Inches(1.5)
            if fig_max_h > max_fig_h:
                fig_max_h = max_fig_h

        # Place the figure
        pic, actual_w, actual_h = _add_image(
            slide, fig["path"], fig_left, fig_top,
            fig_max_w, fig_max_h,
            slide_num=slide_data.get("slide_number", "?"),
        )
        actual_fig_h = actual_h if actual_h else int(fig_max_h)

        # Context/caption text
        context = slide_data.get("context_line")
        caption = slide_data.get("figure_caption") or fig.get("caption", "")

        if fig_on_right:
            # === FIGURE ON RIGHT → ANNOTATIONS ON LEFT ===
            ann_left = int(MARGIN_LEFT)
            ann_width = int(fig_left) - ann_left - Inches(0.3)

            if annotations and ann_width > Inches(1.0):
                _place_annotations_individually(
                    slide, annotations,
                    left=ann_left, top=CONTENT_TOP,
                    width=ann_width, bottom_limit=MAX_CONTENT_BOTTOM,
                )

            # Context/caption below the figure (right-aligned with figure)
            ctx_y = int(fig_top) + actual_fig_h + Inches(0.1)
            if context and ctx_y + Inches(0.35) <= int(MAX_CONTENT_BOTTOM):
                _add_context_line(slide, context, fig_left, ctx_y, fig_max_w)
            elif pic and caption and ctx_y + Inches(0.35) <= int(MAX_CONTENT_BOTTOM):
                _add_caption(slide, caption, fig_left, ctx_y, fig_max_w)

        else:
            # === FIGURE CENTERED → ANNOTATIONS BELOW ===
            fig_bottom = int(fig_top) + actual_fig_h
            cursor_y = fig_bottom + Inches(0.1)

            # Context line or caption
            if context and cursor_y + Inches(0.35) <= int(MAX_CONTENT_BOTTOM):
                _add_context_line(
                    slide, context, MARGIN_LEFT, cursor_y,
                    SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
                )
                cursor_y += Inches(0.45)
            elif pic and caption and cursor_y + Inches(0.40) <= int(MAX_CONTENT_BOTTOM):
                _add_caption(
                    slide, caption, MARGIN_LEFT, cursor_y,
                    SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
                )
                cursor_y += Inches(0.50)

            # Annotations below context/caption — individual textboxes
            if annotations:
                _place_annotations_individually(
                    slide, annotations,
                    left=MARGIN_LEFT, top=cursor_y,
                    width=SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
                    bottom_limit=MAX_CONTENT_BOTTOM,
                )
    else:
        # Fallback: no matching figure — use annotations as text, or bullets
        text_items = annotations or slide_data.get("bullet_points", [])
        if text_items:
            if annotations:
                _place_annotations_individually(
                    slide, text_items,
                    left=MARGIN_LEFT, top=CONTENT_TOP,
                    width=SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
                    bottom_limit=MAX_CONTENT_BOTTOM,
                )
            else:
                _add_bullets(
                    slide, text_items, MARGIN_LEFT, CONTENT_TOP,
                    SLIDE_WIDTH - MARGIN_LEFT - MARGIN_RIGHT,
                )

    _add_references(slide, slide_data.get("references"))
    _add_slide_number(slide, slide_data.get("slide_number", ""))
    _add_speaker_notes(slide, slide_data)
    return slide


def add_key_number_slide(prs, slide_data, figures=None):
    """Key number slide: large centered metric with context and annotations.

    Used for the paper's headline result (e.g. "28.4 BLEU").
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg_white(slide)
    _add_accent_bar(slide)

    _add_title_text(slide, slide_data.get("title", ""))

    cursor = VerticalCursor(start=Inches(2.0))

    # Large key number — centered and prominent
    key_num = slide_data.get("key_number", "")
    if key_num:
        num_h = Inches(2.0)
        num_top = cursor.place(num_h, gap=Inches(0.15))
        txBox = slide.shapes.add_textbox(
            Inches(2.0), num_top,
            SLIDE_WIDTH - Inches(4.0), num_h,
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(key_num)
        run.font.name = FONT_TITLE
        run.font.size = KEY_NUMBER_SIZE
        run.font.bold = True
        run.font.color.rgb = COLOR_ACCENT

    # Context below the number
    context = slide_data.get("key_number_context", "")
    if context:
        ctx_h = Inches(0.6)
        ctx_top = cursor.place(ctx_h, gap=Inches(0.2))
        txBox = slide.shapes.add_textbox(
            Inches(2.0), ctx_top,
            SLIDE_WIDTH - Inches(4.0), ctx_h,
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(context)
        run.font.name = FONT_BODY
        run.font.size = KEY_NUMBER_CONTEXT_SIZE
        run.font.color.rgb = COLOR_BODY

    # Annotations below context — use remaining space
    annotations = slide_data.get("annotations", [])
    if annotations:
        cursor.advance(Inches(0.15))
        ann_h = min(cursor.remaining(int(MAX_CONTENT_BOTTOM)), Inches(1.5))
        if ann_h > Inches(0.4):
            ann_top = cursor.place(ann_h)
            _add_annotations(
                slide, annotations, Inches(2.5), ann_top,
                SLIDE_WIDTH - Inches(5.0), ann_h,
            )

    _add_references(slide, slide_data.get("references"))
    _add_slide_number(slide, slide_data.get("slide_number", ""))
    _add_speaker_notes(slide, slide_data)
    return slide


def add_equation_slide(prs, slide_data, figures=None):
    """Equation slide: assertion title + centered equation + term annotations.

    When figure_reference is present, uses a two-column layout:
      Left: equation + annotations
      Right: mechanism/diagram figure
    Otherwise centers the equation with annotations below.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg_white(slide)
    _add_accent_bar(slide)

    _add_title_text(slide, slide_data.get("title", ""))

    equation = slide_data.get("equation_latex")
    annotations = slide_data.get("annotations", [])
    fig_ref = slide_data.get("figure_reference")
    fig = _find_figure(fig_ref, figures) if figures and fig_ref else None

    if fig:
        # Two-column: equation left, figure right
        eq_left = Inches(0.8)
        eq_width = int(SLIDE_WIDTH * 0.45)
        fig_area_left = int(SLIDE_WIDTH * 0.52)
        fig_area_w = int(SLIDE_WIDTH * 0.42)
        fig_max_h = int(SLIDE_HEIGHT * 0.55)

        cursor = VerticalCursor(start=CONTENT_TOP)

        if equation:
            eq_top = cursor.place(Inches(1.5), gap=Inches(0.2))
            _add_equation(slide, equation, eq_left, eq_top, eq_width)

        if annotations:
            ann_top = cursor.place(Inches(2.0))
            _add_annotations(slide, annotations, eq_left, ann_top, eq_width, Inches(2.0))

        # Figure on the right
        _add_image(
            slide, fig["path"], fig_area_left, int(CONTENT_TOP),
            fig_area_w, fig_max_h,
            slide_num=slide_data.get("slide_number", "?"),
        )
    else:
        # Centered equation
        eq_width = SLIDE_WIDTH - Inches(4.0)
        eq_left = Inches(2.0)

        cursor = VerticalCursor(start=Inches(2.2))

        if equation:
            eq_top = cursor.place(Inches(1.5), gap=Inches(0.3))
            _add_equation(slide, equation, eq_left, eq_top, int(eq_width))

        if annotations:
            ann_remaining = min(cursor.remaining(int(MAX_CONTENT_BOTTOM)), Inches(2.5))
            if ann_remaining > Inches(0.5):
                ann_top = cursor.place(ann_remaining)
                _add_annotations(
                    slide, annotations, Inches(2.0), ann_top,
                    SLIDE_WIDTH - Inches(4.0), ann_remaining,
                )

    # Context line if provided — placed via cursor (shared across both branches)
    context = slide_data.get("context_line")
    if context:
        # For the two-column branch, cursor tracks left-side elements;
        # for centered branch, cursor tracks equation + annotations.
        ctx_top = max(cursor.y, int(Inches(5.8)))
        if ctx_top + Inches(0.35) <= int(MAX_CONTENT_BOTTOM):
            _add_context_line(
                slide, context, Inches(2.0), ctx_top,
                SLIDE_WIDTH - Inches(4.0),
            )

    _add_references(slide, slide_data.get("references"))
    _add_slide_number(slide, slide_data.get("slide_number", ""))
    _add_speaker_notes(slide, slide_data)
    return slide


def add_thankyou_slide(prs, slide_data):
    """Closing Thank You slide with QR code placeholder."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg_white(slide)
    _add_accent_bar(slide)

    _add_title_text(
        slide, "Thank You",
        left=Inches(2.0), top=Inches(2.2),
        width=SLIDE_WIDTH - Inches(4.0),
        size=Pt(44), alignment=PP_ALIGN.CENTER,
    )

    bullets = slide_data.get("bullet_points", [])
    if not bullets:
        bullets = ["Questions?"]

    # Bounded height: 0.4" per bullet, max 1.6"
    bullet_h = min(Inches(len(bullets) * 0.4), Inches(1.6))
    txBox = slide.shapes.add_textbox(
        Inches(2.0), Inches(3.8),
        SLIDE_WIDTH - Inches(4.0), bullet_h,
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = _truncate(bullet)
        run.font.name = FONT_BODY
        run.font.size = Pt(20)
        run.font.color.rgb = COLOR_BODY

    # QR code placeholder — below bullets with gap
    qr_size = Inches(1.2)
    qr_left = (SLIDE_WIDTH - qr_size) // 2
    qr_top = Inches(3.8) + bullet_h + Inches(0.3)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        int(qr_left), int(qr_top), int(qr_size), int(qr_size),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)
    shape.line.color.rgb = COLOR_LIGHT_GRAY
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "QR Code"
    run.font.name = FONT_BODY
    run.font.size = SMALL_SIZE
    run.font.color.rgb = COLOR_LIGHT_GRAY

    _add_slide_number(slide, slide_data.get("slide_number", ""))
    _add_speaker_notes(slide, slide_data)
    return slide


# =============================================================
# LAYOUT DISPATCHER
# =============================================================

LAYOUT_MAP = {
    "title": "title",
    "motivation": "content",
    "insight": "content",
    "research_question": "content",
    "method": "content",
    "method_detail": "content",
    "theory": "table",
    "result": "content",
    "result_table": "table",
    "comparison_table": "table",
    "analysis": "content",
    "visualization": "figure",
    "generalization": "content",
    "comparison": "table",
    "conclusion": "content",
    "thankyou": "thankyou",
    # Survey / literature-review types
    "taxonomy": "content",
    "research_gaps": "content",
    # Theory paper types
    "formulation": "content",
    "proof": "content",
    "implications": "content",
    # Position paper types
    "position": "content",
    "argument": "content",
    # Backup slide types
    "backup": "content",
    "backup_methodology": "content",
    "backup_full_table": "table",
    "backup_ablation": "table",
    "backup_ablation_table": "table",
    "backup_limitations": "content",
    "backup_future_work": "content",
    "backup_reproducibility": "content",
    "backup_extra_experiment": "table",
    "backup_visualization": "figure",
}


def _dispatch_slide(prs, slide_data, figures, paper_meta):
    """Route a single slide_data dict to the correct layout builder.

    Priority order:
    1. title / thankyou (always win)
    2. Assertion-evidence "layout" field (hero_figure, hero_table, key_number, equation)
    3. table_data presence (force table layout)
    4. content_type-based mapping (backward compatibility)
    5. Default: content slide
    """
    if not isinstance(slide_data, dict):
        return  # skip malformed entries

    # Normalize figure_reference — Claude sometimes returns a dict instead of string
    fig_ref = slide_data.get("figure_reference")
    if isinstance(fig_ref, dict):
        num = fig_ref.get("figure_number", "")
        page = fig_ref.get("page_number", fig_ref.get("page", ""))
        slide_data["figure_reference"] = f"Figure {num} (page {page})" if num else None

    content_type = slide_data.get("content_type", "content").lower()
    ct_layout = LAYOUT_MAP.get(content_type, "content")

    # 1. Special layouts that always win
    if ct_layout == "title":
        add_title_slide(prs, slide_data, paper_meta)
        return
    if ct_layout == "thankyou":
        add_thankyou_slide(prs, slide_data)
        return

    # 2. Assertion-evidence "layout" field takes priority
    ae_layout = slide_data.get("layout", "").lower()
    if ae_layout == "key_number" and slide_data.get("key_number"):
        add_key_number_slide(prs, slide_data, figures)
        return
    if ae_layout == "equation" and slide_data.get("equation_latex"):
        add_equation_slide(prs, slide_data, figures)
        return
    if ae_layout == "hero_figure":
        fig = _find_figure(slide_data.get("figure_reference"), figures) if figures else None
        if fig:
            add_figure_slide(prs, slide_data, figures)
            return
    if ae_layout == "hero_table" and slide_data.get("table_data"):
        add_table_slide(prs, slide_data, figures)
        return

    # 3. ANY slide with table_data gets a real table
    if slide_data.get("table_data"):
        add_table_slide(prs, slide_data, figures)
        return

    # 4. content_type-based figure layout (backward compatibility)
    if ct_layout == "figure":
        fig_ref = slide_data.get("figure_reference")
        fig = _find_figure(fig_ref, figures) if figures and fig_ref else None
        if fig:
            add_figure_slide(prs, slide_data, figures)
            return

    # 5. Default: content slide (handles bullets/annotations + optional side figure)
    add_content_slide(prs, slide_data, figures)


def validate_visual_coverage(slide_plan: dict) -> list[str]:
    """Check visual coverage and notes/visual mismatches. Returns warnings."""
    warnings = []

    checklist = slide_plan.get("visual_checklist", {})
    ratio_str = checklist.get("visual_ratio", "0%")
    try:
        ratio = int(ratio_str.replace("%", ""))
    except (ValueError, AttributeError):
        ratio = 0

    if ratio < 40:
        warnings.append(
            f"Low visual ratio ({ratio}%). "
            f"Figures used: {checklist.get('figures_used', [])}. "
            f"Tables used: {checklist.get('tables_used', [])}. "
            f"Skipped: {checklist.get('figures_available_but_skipped', [])} / "
            f"{checklist.get('tables_available_but_skipped', [])}."
        )

    for slide in slide_plan.get("slides", []):
        num = slide.get("slide_number", "?")
        notes = slide.get("speaker_notes", "").lower()
        has_table = slide.get("table_data") is not None
        has_figure = slide.get("figure_reference") is not None

        if "table" in notes and not has_table:
            warnings.append(
                f"Slide {num}: Notes reference a table but no table_data exists."
            )
        if ("figure" in notes or "diagram" in notes) and not has_figure:
            warnings.append(
                f"Slide {num}: Notes reference a figure but no figure_reference exists."
            )

    return warnings


def _normalize_fig_ref(ref) -> str | None:
    """Extract a canonical figure key like 'Figure 3' from a reference string or dict."""
    if isinstance(ref, dict):
        num = ref.get("figure_number", "")
        return f"Figure {num}" if num else None
    if not isinstance(ref, str):
        return None
    m = re.search(r"Figure\s*(\d+)", ref, re.IGNORECASE)
    return f"Figure {m.group(1)}" if m else None


def build_presentation(
    slide_plan: dict,
    figures: list[dict] | None = None,
    output_path: str = "output.pptx",
) -> str:
    """Build a .pptx file from a slide plan.

    Args:
        slide_plan: The slide plan dict from slide_planner.plan_slides().
        figures: List of figure dicts from pdf_parser.parse_pdf()["figures"].
        output_path: Where to save the .pptx file.

    Returns:
        The absolute path to the saved .pptx file.
    """
    figures = figures or []

    # Validate visual coverage (log warnings, don't block)
    warnings = validate_visual_coverage(slide_plan)
    for w in warnings:
        log.warning(w)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    paper_meta = {
        "talk_title": slide_plan.get("talk_title", ""),
        "talk_subtitle": slide_plan.get("talk_subtitle", ""),
        "authors": slide_plan.get("authors", ""),
        "venue": slide_plan.get("venue", ""),
    }

    # --- Main slides ---
    main_slides = slide_plan.get("slides", [])
    used_figures: set[str] = set()
    slide_number = 0  # continuous counter across main + backup

    for slide_data in main_slides:
        if isinstance(slide_data, dict):
            slide_number += 1
            slide_data["slide_number"] = slide_number  # override plan's number
        _dispatch_slide(prs, slide_data, figures, paper_meta)
        key = _normalize_fig_ref(slide_data.get("figure_reference") if isinstance(slide_data, dict) else None)
        if key:
            used_figures.add(key)

    # --- Backup slides (divider + deduplicated) ---
    raw_backups = slide_plan.get("backup_slides", [])
    if raw_backups:
        # Deduplicate: skip backups whose figure is already shown
        filtered_backups: list[dict] = []
        for sd in raw_backups:
            if not isinstance(sd, dict):
                continue
            key = _normalize_fig_ref(sd.get("figure_reference"))
            if key and key in used_figures:
                log.info("Skipping duplicate backup slide: %s (figure %s already used)",
                         sd.get("title", "?"), key)
                continue
            if key:
                used_figures.add(key)
            filtered_backups.append(sd)

        if filtered_backups:
            # Divider slide — no slide number displayed, but occupies a position
            divider = prs.slides.add_slide(prs.slide_layouts[6])
            _set_slide_bg_white(divider)
            _add_accent_bar(divider)
            _add_title_text(
                divider, "Backup Slides",
                left=Inches(2.0), top=Inches(3.0),
                width=SLIDE_WIDTH - Inches(4.0),
                size=Pt(36), alignment=PP_ALIGN.CENTER,
            )
            slide_number += 1  # account for divider occupying a slide position

            for slide_data in filtered_backups:
                if isinstance(slide_data, dict):
                    slide_number += 1
                    slide_data["slide_number"] = slide_number  # continuous numbering
                    # Normalize: backup annotations get bullet markers for consistency
                    anns = slide_data.get("annotations", [])
                    if anns and not slide_data.get("bullet_points"):
                        slide_data["annotations"] = [
                            a if isinstance(a, str) and a.startswith("•") else "•  " + str(a)
                            for a in anns
                        ]
                _dispatch_slide(prs, slide_data, figures, paper_meta)

    # --- Post-processing: fix slide numbers to match actual position ---
    # This corrects any numbering drift regardless of how slides were assembled.
    for i, slide in enumerate(prs.slides):
        is_divider = False
        number_shape = None
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            # Detect the backup divider slide
            if "Backup Slides" in text:
                is_divider = True
            # Detect the slide number textbox: small, bottom-right, numeric
            if (text.isdigit()
                    and shape.width < Inches(1.5)
                    and shape.top > Inches(6.0)):
                number_shape = shape

        if is_divider and number_shape:
            # Clear number on divider slide
            for para in number_shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ""
        elif number_shape:
            # Set correct 1-indexed position
            for para in number_shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = str(i + 1)

    # Save
    out_dir = os.path.dirname(output_path)
    if out_dir:
        os.makedirs(out_dir, exist_ok=True)

    prs.save(output_path)
    return os.path.abspath(output_path)
